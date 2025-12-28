import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from src.utils import get_logger

logger = get_logger(__name__)

class PPTXGenerator:
    def __init__(self, width_inches=13.333, height_inches=7.5):
        self.prs = Presentation()
        self.prs.slide_width = Inches(width_inches)
        self.prs.slide_height = Inches(height_inches)
        self.width_inches = width_inches
        self.height_inches = height_inches

    def _hex_to_rgb(self, hex_color):
        """Converts hex string (e.g., '#FF0000' or '#F00') to RGBColor object."""
        try:
            hex_color = hex_color.lstrip('#')
            if len(hex_color) == 3:
                hex_color = ''.join([c*2 for c in hex_color])
            return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
        except Exception:
            return RGBColor(0, 0, 0) # Default to black

    def add_slide(self, layout_data, bg_image_path, original_width_px, original_height_px, font_family="Malgun Gothic"):
        """
        Adds a new slide to the presentation.
        :param layout_data: List of dicts containing text layout info.
        :param bg_image_path: Path to the background image.
        :param original_width_px: Width of the original image source.
        :param original_height_px: Height of the original image source.
        :param font_family: Font family to use for text.
        """
        # Use blank layout (usually index 6)
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # 1. Add Background Image
        if bg_image_path and os.path.exists(bg_image_path):
            try:
                # Add picture covering the whole slide
                slide.shapes.add_picture(bg_image_path, 0, 0, width=self.prs.slide_width, height=self.prs.slide_height)
            except Exception as e:
                logger.error(f"Failed to add background image to PPTX: {e}")

        # Calculate scaling factors
        scale_x = self.width_inches / original_width_px
        scale_y = self.height_inches / original_height_px

        # 2. Add Text Boxes
        for item in layout_data:
            try:
                bbox = item.get('bbox_px')
                if not bbox: continue

                x_px, y_px, w_px, h_px = bbox
                text = item.get('text', '')
                style = item.get('style', {})

                # Convert to Inches
                left = Inches(x_px * scale_x)
                top = Inches(y_px * scale_y)
                width = Inches(w_px * scale_x)
                height = Inches(h_px * scale_y)

                textbox = slide.shapes.add_textbox(left, top, width, height)
                tf = textbox.text_frame
                tf.word_wrap = True
                
                p = tf.paragraphs[0]
                p.text = text

                # Apply Styling
                # Font Size:
                if 'normalized_font_size_px' in item:
                    font_size_px = item['normalized_font_size_px']
                else:
                    # Fallback calculation if not pre-normalized
                     _, _, _, h_px_item = bbox # We unpacked bbox earlier as x_px, y_px, w_px, h_px
                     # Simple heuristic: 75% of height as font size if single line, but we don't know lines easily.
                     # Let's assume height is close to font size for single line text or use code_generator logic
                     # code_generator: single_line_height = h / line_count. 
                     # We'll use a rough estimate: Height * 0.7
                     # Better: Check checking specific heuristics later if needed.
                     font_size_px = h_px * 0.75
                
                font_size_inches = font_size_px * scale_x
                p.font.size = Pt(font_size_inches * 72)

                # Color
                # Safe access to style
                hex_color = style.get('color', '#000000')
                p.font.color.rgb = self._hex_to_rgb(hex_color)

                # Bold
                if style.get('font_weight') == 'bold':
                    p.font.bold = True
                
                # Alignment
                align = style.get('align', 'left')
                if align == 'center':
                    p.alignment = PP_ALIGN.CENTER
                elif align == 'right':
                    p.alignment = PP_ALIGN.RIGHT
                else:
                    p.alignment = PP_ALIGN.LEFT

                # Font Family
                p.font.name = font_family

            except Exception as e:
                logger.error(f"Error adding text to PPTX: {e}")

    def save(self, output_path):
        self.prs.save(output_path)
        logger.info(f"PPTX saved to {output_path}")
